import json
import base64
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE
import requests
from PIL import Image
import os
import re

# Base slide sizes
BASE_SIZES = {
    '720p': (1280, 720),
    '1080p': (1920, 1080),
    '1440p': (2560, 1440),
    '4K': (3840, 2160),
}

def safe_int_conversion(value, default=0):
    if value is None or value == 'auto' or value == 'normal':
        return default
    try:
        if isinstance(value, str):
            numeric_part = re.search(r'[\d.]+', value)
            if numeric_part:
                return int(float(numeric_part.group()))
        return int(float(value))
    except (ValueError, TypeError):
        return default

def safe_float_conversion(value, default=0.0):
    if value is None or value == 'auto' or value == 'normal':
        return default
    try:
        if isinstance(value, str):
            numeric_part = re.search(r'[\d.]+', value)
            if numeric_part:
                return float(numeric_part.group())
        return float(value)
    except (ValueError, TypeError):
        return default

def analyze_content_bounds(slides_data):
    """Analyze all slides to determine the exact content bounds"""
    max_right = 0
    max_bottom = 0
    
    print("Analyzing content bounds...")
    
    for slide_info in slides_data:
        for element in slide_info.get('elements', []):
            x = element.get('x', 0)
            y = element.get('y', 0)
            width = element.get('width', 0)
            height = element.get('height', 0)
            
            right = x + width
            bottom = y + height
            
            max_right = max(max_right, right)
            max_bottom = max(max_bottom, bottom)
    
    print(f"Content bounds: {max_right} x {max_bottom} pixels")
    return max_right, max_bottom

def calculate_optimal_slide_size(content_width, content_height, base_size='1080p', padding=20):
    """Calculate slide size based on content dimensions with specified padding."""
    base_width, base_height = BASE_SIZES.get(base_size, BASE_SIZES['1080p'])
    
    if content_width > 0 and content_height > 0:
        final_width = content_width + (padding * 2)
        final_height = content_height + (padding * 2)
        print(f"Using content dimensions {content_width}x{content_height} with {padding}px padding all sides")
    else:
        final_width = base_width + (padding * 2)
        final_height = base_height + (padding * 2)
        print(f"No content dimensions specified, using base size {base_size} with {padding}px padding")
    
    min_width = 800
    min_height = 600
    
    if final_width < min_width:
        final_width = min_width
        print(f"Using minimum width: {min_width}")
    
    if final_height < min_height:
        final_height = min_height
        print(f"Using minimum height: {min_height}")
    
    final_width = int(final_width)
    final_height = int(final_height)
    
    size_name = f"Custom {final_width}x{final_height}"
    if final_width == base_width + (padding * 2) and final_height == base_height + (padding * 2):
        size_name = base_size
    
    print(f"Selected slide size: {size_name} ({final_width}x{final_height})")
    return final_width, final_height, size_name

def center_content_on_slide(slides_data, slide_width, slide_height, content_width, content_height, padding=20):
    """Center content on the slide, ensuring at least the specified padding on all sides."""
    offset_x = max(padding, (slide_width - content_width - padding * 2) // 2 + padding)
    offset_y = max(padding, (slide_height - content_height - padding * 2) // 2 + padding)
    
    if offset_x > padding or offset_y > padding:
        print(f"Centering content with offset: ({offset_x}, {offset_y})")
    
    for slide_info in slides_data:
        for element in slide_info.get('elements', []):
            element['x'] = element.get('x', 0) + offset_x
            element['y'] = element.get('y', 0) + offset_y
    
    return slides_data

def parse_color(color_str):
    if not color_str or color_str in ['transparent', 'rgba(0, 0, 0, 0)']:
        return None
    
    if color_str.startswith('rgb'):
        color_values = re.findall(r'\d+', color_str)
        if len(color_values) >= 3:
            return RGBColor(int(color_values[0]), int(color_values[1]), int(color_values[2]))
    
    elif color_str.startswith('#'):
        color_str = color_str.lstrip('#')
        if len(color_str) == 6:
            return RGBColor(int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16))
        elif len(color_str) == 3:
            return RGBColor(int(color_str[0]*2, 16), int(color_str[1]*2, 16), int(color_str[2]*2, 16))
    
    named_colors = {
        'black': RGBColor(0, 0, 0),
        'white': RGBColor(255, 255, 255),
        'red': RGBColor(255, 0, 0),
        'green': RGBColor(0, 128, 0),
        'blue': RGBColor(0, 0, 255),
        'yellow': RGBColor(255, 255, 0),
        'gray': RGBColor(128, 128, 128),
        'grey': RGBColor(128, 128, 128),
    }
    
    return named_colors.get(color_str.lower())

def parse_border(styles):
    if not styles.get('border') or styles['border'] == 'none':
        return None, None, None

    border_width = safe_float_conversion(styles.get('borderWidth', '0').replace('px', ''))
    border_style = styles.get('borderStyle', 'solid')
    border_color = parse_color(styles.get('borderColor'))

    if not (border_width > 0 and border_style and border_color):
        border_str = styles.get('border', '')
        parts = border_str.split()
        for part in parts:
            if 'px' in part:
                border_width = safe_float_conversion(part.replace('px', ''))
            elif part in ['solid', 'dashed', 'dotted']:
                border_style = part
            elif part.startswith('rgb') or part.startswith('#'):
                border_color = parse_color(part)

        if not (border_width > 0 and border_style and border_color):
            match = re.match(r'(\d*\.?\d*)px\s+(\w+)\s+(.+)', border_str)
            if match:
                border_width = safe_float_conversion(match.group(1))
                border_style = match.group(2)
                border_color = parse_color(match.group(3))

    if border_width > 0:
        border_width = max(0.5, border_width * 0.75)
    
    return border_width, border_style, border_color

def parse_border_radius(radius_str, shape_width, shape_height):
    if not radius_str or radius_str == '0px':
        return 0
    
    try:
        match = re.search(r'[\d.]+', radius_str)
        if match:
            radius_px = float(match.group())
            min_dimension = min(shape_width, shape_height)
            radius_ratio = min(radius_px / min_dimension, 0.5)
            return radius_ratio
    except (ValueError, TypeError):
        return 0

def pixels_to_emu(pixels):
    return int(pixels * 9525)

def get_font_size_pt(font_size_px):
    if font_size_px <= 0:
        return 12
    return max(8, int(font_size_px * 0.8))

def add_separator_element(slide, element, slide_width, slide_height, debug=False):
    x, y, width = element['x'], element['y'], element['width']
    styles = element.get('styles', {})
    
    border_width, border_style, border_color = parse_border(styles)
    if not (border_width and border_style and border_color):
        border_width = safe_float_conversion(styles.get('borderTopWidth', '1').replace('px', ''), 1)
        border_style = styles.get('borderTopStyle', 'solid')
        border_color = parse_color(styles.get('borderTopColor', '#E7EAE8'))
    
    if not border_width or not border_color:
        return

    x = max(0, min(x, slide_width - 1))
    y = max(0, min(y, slide_height - 1))
    width = max(1, min(width, slide_width - x))
    
    try:
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            pixels_to_emu(x),
            pixels_to_emu(y),
            pixels_to_emu(x + width),
            pixels_to_emu(y)
        )
        
        line = connector.line
        line.width = Pt(max(0.5, border_width * 0.75))
        line.color.rgb = border_color
        
        if border_style == 'dashed':
            line.dash_style = MSO_LINE.DASH
        elif border_style == 'dotted':
            line.dash_style = MSO_LINE.ROUND_DOT
        
        if debug:
            print(f"Added separator line at ({x}, {y}) width {width}, border: {border_width}px {border_style} {border_color}")
    
    except Exception as e:
        print(f"Failed to add separator element: {e}")

def add_text_element(slide, element, slide_width, slide_height, debug=False):
    x, y, width, height = element['x'], element['y'], element['width'], element['height']
    text = element.get('text', '').strip()
    
    if not text:
        return

    x = max(0, min(x, slide_width - 10))
    y = max(0, min(y, slide_height - 10))
    width = max(10, min(width, slide_width - x))
    height = max(10, min(height, slide_height - y))
    
    try:
        textbox = slide.shapes.add_textbox(
            pixels_to_emu(x), pixels_to_emu(y),
            pixels_to_emu(width), pixels_to_emu(height)
        )
        
        # shadow textbox
        textbox.shadow.inherit = False
        
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.auto_size = None
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        
        paragraph = text_frame.paragraphs[0]
        styles = element.get('styles', {})

        font_size_px = safe_float_conversion(styles.get('fontSize', '12').replace('px', ''))
        if font_size_px > 0:
            font_size_pt = get_font_size_pt(font_size_px)
            paragraph.font.size = Pt(font_size_pt)
            if debug:
                print(f"Font size: {font_size_px}px -> {font_size_pt}pt")
        
        font_family = styles.get('fontFamily', 'Arial')
        if font_family:
            if ',' in font_family:
                font_family = font_family.split(',')[0].strip()
            paragraph.font.name = font_family
        
        font_weight = styles.get('fontWeight', 'normal')
        if font_weight == 'bold' or safe_int_conversion(font_weight) >= 700:
            paragraph.font.bold = True

        if styles.get('fontStyle') == 'italic':
            paragraph.font.italic = True

        text_color = parse_color(styles.get('color', 'black'))
        if text_color:
            paragraph.font.color.rgb = text_color
        
        # textshadow dis
        paragraph.font.shadow = False
        
        text_align = styles.get('textAlign', 'left')
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        paragraph.alignment = alignment_map.get(text_align, PP_ALIGN.LEFT)
        
        bg_color = parse_color(styles.get('backgroundColor'))
        if bg_color:
            fill = textbox.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        else:
            textbox.fill.background()

        textbox.line.fill.background()
        
        padding_top = safe_int_conversion(styles.get('paddingTop', '0').replace('px', ''))
        padding_right = safe_int_conversion(styles.get('paddingRight', '0').replace('px', ''))
        padding_bottom = safe_int_conversion(styles.get('paddingBottom', '0').replace('px', ''))
        padding_left = safe_int_conversion(styles.get('paddingLeft', '0').replace('px', ''))
        
        if padding_top > 0:
            text_frame.margin_top = pixels_to_emu(padding_top)
        if padding_right > 0:
            text_frame.margin_right = pixels_to_emu(padding_right)
        if padding_bottom > 0:
            text_frame.margin_bottom = pixels_to_emu(padding_bottom)
        if padding_left > 0:
            text_frame.margin_left = pixels_to_emu(padding_left)
        
        if debug:
            print(f"Added text: '{text[:50]}...' at ({x}, {y}) size ({width}x{height})")
    
    except Exception as e:
        print(f"Failed to add text element: {e}")

def add_shape_element(slide, element, slide_width, slide_height, debug=False):
    x, y, width, height = element['x'], element['y'], element['width'], element['height']
    styles = element.get('styles', {})
    text = element.get('text', '').strip()
    
    bg_color = parse_color(styles.get('backgroundColor'))
    border_width, border_style, border_color = parse_border(styles)
    border_radius = parse_border_radius(styles.get('borderRadius', '0px'), width, height)

    if not bg_color and not border_width and not text:
        return

    x = max(0, min(x, slide_width - 1))
    y = max(0, min(y, slide_height - 1))
    width = max(1, min(width, slide_width - x))
    height = max(1, min(height, slide_height - y))
    
    try:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if border_radius > 0 else MSO_SHAPE.RECTANGLE
        
        shape = slide.shapes.add_shape(
            shape_type,
            pixels_to_emu(x), pixels_to_emu(y),
            pixels_to_emu(width), pixels_to_emu(height)
        )
        
        # shadow disablw
        shape.shadow.inherit = False

        if border_radius > 0 and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
            try:
                shape.adjustments[0] = border_radius
            except:
                pass
        
        if bg_color:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        else:
            shape.fill.background()

        if border_width and border_width > 0:
            line = shape.line
            line.width = Pt(border_width)
            if border_color:
                line.color.rgb = border_color
            if border_style == 'dashed':
                line.dash_style = MSO_LINE.DASH
            else:
                line.dash_style = MSO_LINE.SOLID
        else:
            shape.line.fill.background()

        if text:
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.text = text
            text_frame.word_wrap = True
            text_frame.auto_size = None
            text_frame.vertical_anchor = MSO_ANCHOR.TOP

            text_frame.margin_left = pixels_to_emu(safe_int_conversion(styles.get('paddingLeft', '0').replace('px', '')))
            text_frame.margin_right = pixels_to_emu(safe_int_conversion(styles.get('paddingRight', '0').replace('px', '')))
            text_frame.margin_top = pixels_to_emu(safe_int_conversion(styles.get('paddingTop', '0').replace('px', '')))
            text_frame.margin_bottom = pixels_to_emu(safe_int_conversion(styles.get('paddingBottom', '0').replace('px', '')))
            
            paragraph = text_frame.paragraphs[0]
            
            font_size_px = safe_float_conversion(styles.get('fontSize', '12').replace('px', ''))
            if font_size_px > 0:
                font_size_pt = get_font_size_pt(font_size_px)
                paragraph.font.size = Pt(font_size_pt)
            
            font_family = styles.get('fontFamily', 'Arial')
            if font_family:
                if ',' in font_family:
                    font_family = font_family.split(',')[0].strip()
                paragraph.font.name = font_family
            
            font_weight = styles.get('fontWeight', 'normal')
            if font_weight == 'bold' or safe_int_conversion(font_weight) >= 700:
                paragraph.font.bold = True

            if styles.get('fontStyle') == 'italic':
                paragraph.font.italic = True

            text_color = parse_color(styles.get('color', 'black'))
            if text_color:
                paragraph.font.color.rgb = text_color
            
            # Disable shadow for text in shape
            paragraph.font.shadow = False
            
            text_align = styles.get('textAlign', 'left')
            alignment_map = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT,
                'justify': PP_ALIGN.JUSTIFY
            }
            paragraph.alignment = alignment_map.get(text_align, PP_ALIGN.LEFT)
        
        if debug:
            print(f"Added shape at ({x}, {y}) size ({width}x{height}), radius={border_radius}, text='{text[:50]}...'")
    
    except Exception as e:
        print(f"Failed to add shape element: {e}")

def add_image_element(slide, element, slide_width, slide_height, debug=False):
    x, y, width, height = element['x'], element['y'], element['width'], element['height']
    img_src = element.get('src', '')
    
    if not img_src:
        return

    x = max(0, min(x, slide_width - 1))
    y = max(0, min(y, slide_height - 1))
    width = max(1, min(width, slide_width - x))
    height = max(1, min(height, slide_height - y))
    
    try:
        temp_path = None
        if img_src.startswith('data:'):
            header, data = img_src.split(',', 1)
            img_data = base64.b64decode(data)
            temp_path = 'temp_image.png'
            with open(temp_path, 'wb') as f:
                f.write(img_data)
        elif img_src.startswith('http'):
            response = requests.get(img_src, timeout=10)
            if response.status_code == 200:
                temp_path = 'temp_image.png'
                with open(temp_path, 'wb') as f:
                    f.write(response.content)
            else:
                print(f"Failed to download image: {img_src}, status: {response.status_code}")
                return
        else:
            if os.path.exists(img_src):
                temp_path = img_src
            else:
                print(f"Image file not found: {img_src}")
                return

        try:
            with Image.open(temp_path) as img:
                img.verify()
        except Exception as e:
            print(f"Invalid image file: {temp_path}, error: {e}")
            if temp_path != img_src and os.path.exists(temp_path):
                os.remove(temp_path)
            return
        
        picture = slide.shapes.add_picture(
            temp_path,
            pixels_to_emu(x), pixels_to_emu(y),
            pixels_to_emu(width), pixels_to_emu(height)
        )
        
        # Disable shadow for the image
        picture.shadow.inherit = False

        if temp_path != img_src and os.path.exists(temp_path):
            os.remove(temp_path)
        
        if debug:
            print(f"Added image: {img_src} at ({x}, {y}) size ({width}x{height})")
    
    except Exception as e:
        print(f"Failed to add image element: {e}")
        if temp_path and temp_path != img_src and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except:
                pass

def create_pptx_from_json(json_path, output_path=None, debug=False, base_size='1080p', padding=20, center_content=True):
    """Create PowerPoint presentation from JSON with HTML-like content fitting."""
    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
    except Exception as e:
        print(f"Error reading JSON file: {e}")
        return

    print(f"Loaded {len(slides_data)} slides from {json_path}")
    
    content_width, content_height = analyze_content_bounds(slides_data)
    slide_width, slide_height, size_name = calculate_optimal_slide_size(content_width, content_height, base_size, padding)
    
    if center_content:
        slides_data = center_content_on_slide(slides_data, slide_width, slide_height, content_width, content_height, padding)
    
    prs = Presentation()
    prs.slide_width = pixels_to_emu(slide_width)
    prs.slide_height = pixels_to_emu(slide_height)
    
    print(f"Creating {size_name} presentation")
    print(f"Slide dimensions: {slide_width}x{slide_height} pixels")
    print(f"Content dimensions: {content_width}x{content_height} pixels")
    
    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        elements = slide_info.get('elements', [])
        
        if debug:
            print(f"\nProcessing slide {slide_info.get('slideId', 'Unknown')} with {len(elements)} elements")

        elements_sorted = sorted(elements, key=lambda e: e.get('zIndex', 0))

        for element in elements_sorted:
            element_type = element.get('type', '').lower()
            class_name = element.get('className', '')
            
            if debug:
                print(f"Processing {element_type} at ({element['x']}, {element['y']}) size ({element['width']}x{element['height']})")

            if element_type == 'img':
                add_image_element(slide, element, slide_width, slide_height, debug)
            
            elif element_type == 'div':
                styles = element.get('styles', {})
                has_background = styles.get('backgroundColor') and styles['backgroundColor'] != 'rgba(0, 0, 0, 0)'
                has_border = styles.get('border') and styles['border'] != 'none'
                has_border_radius = styles.get('borderRadius') and styles['borderRadius'] != '0px'

                if 'separator' in class_name:
                    add_separator_element(slide, element, slide_width, slide_height, debug)
                elif has_background or has_border or has_border_radius:
                    add_shape_element(slide, element, slide_width, slide_height, debug)
                elif element.get('text'):
                    add_text_element(slide, element, slide_width, slide_height, debug)
            
            elif element_type in ['span', 'p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'] and element.get('text'):
                add_text_element(slide, element, slide_width, slide_height, debug)
    
    if output_path is None:
        base_name = os.path.splitext(os.path.basename(json_path))[0]
        clean_size_name = size_name.replace(' ', '_').replace('(', '').replace(')', '').replace('x', 'x')
        output_path = f"{base_name}_{clean_size_name}.pptx"
    
    try:
        prs.save(output_path)
        print(f"\nPowerPoint presentation saved successfully as '{output_path}'")
        print(f"Final slide size: {size_name}")
        
        base_w, base_h = BASE_SIZES[base_size]
        if slide_width > base_w and slide_height > base_h:
            print(f"✓ Extended both width and height from {base_size} to fit content")
        elif slide_width > base_w:
            print(f"✓ Extended width from {base_size} to fit content")
        elif slide_height > base_h:
            print(f"✓ Extended height from {base_size} to fit content")
        else:
            print(f"✓ Content fits within {base_size} dimensions")
            
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == "__main__":
    create_pptx_from_json('slides_data.json', 'output_dynamic.pptx', debug=True, padding=50)