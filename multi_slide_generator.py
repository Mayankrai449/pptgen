import json
import base64
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE, MSO_COLOR_TYPE
import requests
from PIL import Image, ImageDraw
import os
import re
import math
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

def safe_int(value, default=0):
    try:
        return int(float(value))
    except (ValueError, TypeError):
        return default

def safe_float(value, default=0.0):
    """Safely convert value to float with better error handling"""
    try:
        if isinstance(value, str):
            # Remove units and extract numeric value
            value = re.sub(r'[^\d.-]', '', value)
        return float(value) if value else default
    except (ValueError, TypeError):
        return default

def pixels_to_emu(pixels):
    """Convert pixels to EMU with high precision"""
    return int(round(pixels * 9525))

def px_to_pt(px):
    """Convert pixels to points"""
    return px * 0.75

def get_font_size_pt(font_size_px):
    if font_size_px <= 0:
        return 12
    return max(6, int(px_to_pt(font_size_px)))

def parse_border_radius(radius_str, shape_width_px, shape_height_px):
    """Parse border radius from CSS string"""
    if not radius_str or radius_str == '0px':
        return 0
    try:
        if isinstance(radius_str, str):
            if '%' in radius_str:
                percent = safe_float(radius_str.replace('%', ''))
                return min(percent / 100, 0.5)
            match = re.search(r'[\d.]+', radius_str)
            if match:
                radius_px = float(match.group())
            else:
                return 0
        else:
            radius_px = float(radius_str)
        min_dimension = min(shape_width_px, shape_height_px)
        if min_dimension > 0:
            return min(radius_px / min_dimension, 0.5)
        return 0
    except (ValueError, TypeError):
        return 0

def parse_color(color_str):
    """Enhanced color parsing with better RGB extraction"""
    if not color_str or color_str in ['transparent', 'rgba(0, 0, 0, 0)', 'none', 'initial', 'inherit']:
        return None
    
    # Handle RGB/RGBA
    if color_str.startswith(('rgb', 'rgba')):
        parts = re.findall(r'[\d.]+', color_str)
        if len(parts) >= 3:
            r = min(255, max(0, int(float(parts[0]))))
            g = min(255, max(0, int(float(parts[1]))))
            b = min(255, max(0, int(float(parts[2]))))
            return RGBColor(r, g, b)
    
    # Handle hex colors
    elif color_str.startswith('#'):
        color_str = color_str.lstrip('#')
        if len(color_str) == 6:
            try:
                return RGBColor(int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16))
            except ValueError:
                pass
        elif len(color_str) == 3:
            try:
                return RGBColor(int(color_str[0]*2, 16), int(color_str[1]*2, 16), int(color_str[2]*2, 16))
            except ValueError:
                pass
    
    # Named colors
    named_colors = {
        'black': RGBColor(0, 0, 0), 'white': RGBColor(255, 255, 255),
        'red': RGBColor(255, 0, 0), 'green': RGBColor(0, 128, 0),
        'blue': RGBColor(0, 0, 255), 'yellow': RGBColor(255, 255, 0),
        'gray': RGBColor(128, 128, 128), 'grey': RGBColor(128, 128, 128),
        'silver': RGBColor(192, 192, 192), 'maroon': RGBColor(128, 0, 0),
        'olive': RGBColor(128, 128, 0), 'lime': RGBColor(0, 255, 0),
        'aqua': RGBColor(0, 255, 255), 'teal': RGBColor(0, 128, 128),
        'navy': RGBColor(0, 0, 128), 'fuchsia': RGBColor(255, 0, 255),
        'purple': RGBColor(128, 0, 128)
    }
    
    return named_colors.get(color_str.lower())

def is_uniform_border(styles):
    widths = [safe_float(styles.get(f'border{side}Width', '0px')) for side in ['Top', 'Right', 'Bottom', 'Left']]
    styles_list = [styles.get(f'border{side}Style', 'none') for side in ['Top', 'Right', 'Bottom', 'Left']]
    colors = [styles.get(f'border{side}Color', '') for side in ['Top', 'Right', 'Bottom', 'Left']]
    return len(set(widths)) == 1 and len(set(styles_list)) == 1 and len(set(colors)) == 1 and widths[0] > 0

def has_any_border(styles):
    """Enhanced border detection"""
    for side in ['Top', 'Right', 'Bottom', 'Left']:
        width_key = f'border{side}Width'
        style_key = f'border{side}Style'
        color_key = f'border{side}Color'
        
        width = safe_float(styles.get(width_key, '0px'))
        style = styles.get(style_key, 'none')
        color = parse_color(styles.get(color_key, ''))
        
        if width > 0 and style not in ['none', 'hidden'] and color is not None:
            return True
    return False

def get_border_info(styles):
    """Enhanced border information extraction"""
    border_info = {}
    for side in ['Top', 'Right', 'Bottom', 'Left']:
        width_key = f'border{side}Width'
        style_key = f'border{side}Style'
        color_key = f'border{side}Color'
        
        width = safe_float(styles.get(width_key, '0px'))
        style = styles.get(style_key, 'none')
        color = parse_color(styles.get(color_key, ''))
        
        border_info[side.lower()] = {
            'width': width,
            'style': style,
            'color': color,
            'has_border': width > 0 and style not in ['none', 'hidden'] and color is not None
        }
    return border_info

def make_rounded_image(input_path, output_path, radius):
    im = Image.open(input_path).convert("RGBA")
    mask = Image.new("L", im.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle((0, 0) + im.size, radius=radius, fill=255)
    im.putalpha(mask)
    im.save(output_path, "PNG")

def create_precise_border_shapes(slide, x, y, width, height, border_info, border_radius=0):
    """Create precise border shapes with accurate positioning"""
    shapes_created = []
    # Always check and draw all four borders if needed
    for side, info in border_info.items():
        if not info['has_border']:
            continue
        border_width = info['width']
        border_color = info['color']
        # Calculate precise line positions
        if side == 'top':
            line_x = x
            line_y = y
            line_width = width
            line_height = border_width
        elif side == 'right':
            line_x = x + width - border_width
            line_y = y
            line_width = border_width
            line_height = height
        elif side == 'bottom':
            line_x = x
            line_y = y + height - border_width
            line_width = width
            line_height = border_width
        elif side == 'left':
            line_x = x
            line_y = y
            line_width = border_width
            line_height = height
        try:
            # Always use rectangle for side borders, rounded only for full shape
            shape_type = MSO_SHAPE.RECTANGLE
            border_shape = slide.shapes.add_shape(
                shape_type,
                pixels_to_emu(line_x), pixels_to_emu(line_y),
                pixels_to_emu(line_width), pixels_to_emu(line_height)
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = border_color
            border_shape.line.fill.background()
            border_shape.shadow.inherit = False
            shapes_created.append(border_shape)
        except Exception as e:
            print(f"Error creating {side} border: {e}")
    return shapes_created

def add_bg_shape(slide, styles, x, y, width, height):
    """Enhanced background shape creation with precise positioning"""
    bg_color = parse_color(styles.get('backgroundColor'))
    border_radius_str = styles.get('borderRadius', '0px')
    border_radius = parse_border_radius(border_radius_str, width, height)
    has_radius = border_radius > 0
    box_shadow = styles.get('boxShadow', 'none')
    has_shadow = box_shadow != 'none'
    has_uniform_border = is_uniform_border(styles)
    has_any_border_sides = has_any_border(styles)
    shapes_created = []
    # Create main background shape
    if bg_color or has_uniform_border or has_radius or has_shadow:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if has_radius else MSO_SHAPE.RECTANGLE
        try:
            bg_shape = slide.shapes.add_shape(
                shape_type,
                pixels_to_emu(x), pixels_to_emu(y),
                pixels_to_emu(width), pixels_to_emu(height)
            )
            if has_radius:
                bg_shape.adjustments[0] = border_radius
            if bg_color:
                bg_shape.fill.solid()
                bg_shape.fill.fore_color.rgb = bg_color
            else:
                bg_shape.fill.background()
            bg_shape.shadow.inherit = False
            # Handle uniform borders
            if has_uniform_border:
                border_width = safe_float(styles.get('borderTopWidth', '1px'))
                border_color = parse_color(styles.get('borderTopColor', 'black'))
                border_style = styles.get('borderTopStyle', 'solid')
                if border_color:
                    bg_shape.line.width = Pt(max(0.5, border_width))
                    bg_shape.line.color.rgb = border_color
                    if border_style == 'dashed':
                        bg_shape.line.dash_style = MSO_LINE.DASH
                    elif border_style == 'dotted':
                        bg_shape.line.dash_style = MSO_LINE.ROUND_DOT
                else:
                    bg_shape.line.fill.background()
            else:
                bg_shape.line.fill.background()
            if has_shadow:
                apply_shadow(bg_shape, box_shadow)
            shapes_created.append(bg_shape)
        except Exception as e:
            print(f"Error adding bg shape: {e}")
    # Always handle non-uniform borders for all sides
    if has_any_border_sides and not has_uniform_border:
        border_info = get_border_info(styles)
        border_shapes = create_precise_border_shapes(
            slide, x, y, width, height, border_info,
            border_radius * min(width, height) if has_radius else 0
        )
        shapes_created.extend(border_shapes)
    return shapes_created

def apply_shadow(shape, box_shadow_str):
    if box_shadow_str == 'none':
        return
    parts = box_shadow_str.split()
    if len(parts) < 3:
        return
    offset_x_px = safe_float(parts[0])
    offset_y_px = safe_float(parts[1])
    blur_px = safe_float(parts[2])
    spread_px = safe_float(parts[3]) if len(parts) > 3 else 0
    color_str = parts[4] if len(parts) > 4 else parts[3]
    if offset_x_px == 0 and offset_y_px == 0 and blur_px == 0 and spread_px == 0:
        return
    color = parse_color(color_str)
    alpha = 1.0
    if 'rgba' in color_str:
        color_parts = re.findall(r'\d+', color_str)
        if len(color_parts) == 4:
            alpha = float(color_parts[3]) / 255
    if not color:
        return
    distance_px = math.sqrt(offset_x_px**2 + offset_y_px**2)
    direction = math.degrees(math.atan2(offset_y_px, offset_x_px)) if distance_px > 0 else 0
    shape.shadow.inherit = False
    shape.shadow.blur = Pt(px_to_pt(blur_px))
    shape.shadow.distance = Pt(px_to_pt(distance_px))
    shape.shadow.angle = direction
    shape.shadow.color.type = MSO_COLOR_TYPE.RGB
    shape.shadow.color.rgb = color
    shape.shadow.transparency = 1 - alpha

def add_inline_group_element(slide, element, slide_width, slide_height, parent_has_shadow=False):
    inline_group = element.get('inlineGroup')
    if not inline_group:
        return
    inline_elements = inline_group.get('inlineElements', [])
    if not inline_elements:
        return
    has_content = any(elem.get('text', '').strip() for elem in inline_elements)
    if not has_content:
        return
    group_rect = inline_group.get('groupRect', {})
    x = max(0, min(group_rect.get('x', 0), slide_width - 10))
    y = max(0, min(group_rect.get('y', 0), slide_height - 10))
    width = max(10, min(group_rect.get('width', 100), slide_width - x))
    height = max(10, min(group_rect.get('height', 20), slide_height - y))
    styles = inline_group.get('styles', {})
    box_shadow = styles.get('boxShadow', 'none')
    has_shadow = box_shadow != 'none' and not parent_has_shadow
    bg_color = parse_color(styles.get('backgroundColor'))
    border_radius_str = styles.get('borderRadius', '0px')
    border_radius = parse_border_radius(border_radius_str, width, height)
    has_radius = border_radius > 0
    has_border = is_uniform_border(styles)
    has_any_border_sides = has_any_border(styles)
    try:
        if bg_color or has_border or has_any_border_sides or has_radius or has_shadow:
            add_bg_shape(slide, styles, x, y, width, height)
        textbox = slide.shapes.add_textbox(
            pixels_to_emu(x), pixels_to_emu(y),
            pixels_to_emu(width), pixels_to_emu(height)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.margin_left = pixels_to_emu(safe_float(styles.get('paddingLeft', '0px').replace('px', '')))
        text_frame.margin_right = pixels_to_emu(safe_float(styles.get('paddingRight', '0px').replace('px', '')))
        text_frame.margin_top = pixels_to_emu(safe_float(styles.get('paddingTop', '0px').replace('px', '')))
        text_frame.margin_bottom = pixels_to_emu(safe_float(styles.get('paddingBottom', '0px').replace('px', '')))
        textbox.fill.background()
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        text_frame.clear()
        p = text_frame.paragraphs[0]  # Use the first paragraph instead of adding new one
        first = True
        text_align = styles.get('textAlign', 'left')
        alignment = PP_ALIGN.CENTER if text_align == 'center' else PP_ALIGN.RIGHT if text_align == 'right' else PP_ALIGN.LEFT
        p.alignment = alignment
        for inline_element in inline_elements:
            if inline_element.get('type') == 'br':
                if p is not None:
                    p = text_frame.add_paragraph()
                    p.alignment = alignment
                first = True
                continue
            element_text = inline_element.get('text', '')
            if first:
                element_text = element_text.lstrip()
            if not element_text.strip():
                continue
            first = False
            run = p.add_run()
            run.text = element_text
            inline_styles = inline_element.get('styles', {})
            font = run.font
            font_size_px = safe_float(inline_styles.get('fontSize', '16').replace('px', ''))
            font.name = inline_styles.get('fontFamily', 'Segoe UI').split(',')[0].strip('"\'')
            font.size = Pt(get_font_size_pt(font_size_px))
            font.bold = inline_styles.get('fontWeight', '400') in ['bold', '600', '700', '800', '900']
            font.italic = inline_styles.get('fontStyle', 'normal') == 'italic'
            color = parse_color(inline_styles.get('color'))
            if color:
                font.color.rgb = color
        # Trim trailing spaces from the last run in the last paragraph
        if text_frame.paragraphs and text_frame.paragraphs[-1].runs:
            last_run = text_frame.paragraphs[-1].runs[-1]
            last_run.text = last_run.text.rstrip()
    except Exception as e:
        print(f"Failed to add inline group element: {e}")

def add_list_paragraphs(text_frame, list_info, level=0, counters=None):
    if counters is None:
        counters = {}
    list_type = list_info.get('type')
    is_ordered = list_type == 'ol'
    list_styles = list_info.get('listStyles', {})
    list_style_type = list_styles.get('listStyleType', 'disc' if not is_ordered else 'decimal')
    if is_ordered:
        counter_key = f'ol_{level}'
        counters[counter_key] = list_info.get('start', 1) - 1
    # Calculate average space_after
    items = list_info.get('items', [])
    avg_space_px = 0
    if len(items) > 1:
        spaces = []
        for i in range(len(items) - 1):
            item_bottom = items[i]['rect']['y'] + items[i]['rect']['height']
            next_top = items[i+1]['rect']['y']
            space_px = next_top - item_bottom
            if space_px > 0:
                spaces.append(space_px)
        avg_space_px = sum(spaces) / len(spaces) if spaces else 0
    space_after_pt = px_to_pt(avg_space_px)
    # Line spacing
    line_height_str = list_info.get('styles', {}).get('lineHeight', 'normal')
    if line_height_str == 'normal':
        line_spacing = 1.15
    else:
        try:
            line_spacing = float(line_height_str)
        except ValueError:
            line_spacing = 1.15
    
    first_item = True
    for item in items:
        p = None
        item_styles = item.get('styles', {})
        default_font_size_px = safe_float(item_styles.get('fontSize', '16').replace('px', ''))
        default_font_size_pt = get_font_size_pt(default_font_size_px)
        default_font_name = item_styles.get('fontFamily', 'Segoe UI').split(',')[0].strip('"\'')
        default_color = parse_color(item_styles.get('color'))
        if item.get('inlineGroup') and item['inlineGroup'].get('inlineElements'):
            inline_elements = item['inlineGroup']['inlineElements']
            first = True
            bullet_added = False
            for inline_element in inline_elements:
                if inline_element.get('type') == 'br':
                    if p is not None:
                        p = text_frame.add_paragraph()
                        p.level = level
                        p.space_after = Pt(space_after_pt)
                        p.line_spacing = line_spacing
                        indent_pt = 18 * level
                        p.left_indent = Pt(indent_pt)
                        p.first_line_indent = Pt(-18)
                    first = True
                    continue
                element_text = inline_element.get('text', '')
                if p is None:
                    if first_item and level == 0:
                        p = text_frame.paragraphs[0]  # Use first paragraph for first item
                    else:
                        p = text_frame.add_paragraph()
                    p.level = level
                    p.space_after = Pt(space_after_pt)
                    p.line_spacing = line_spacing
                    indent_pt = 18 * level
                    p.left_indent = Pt(indent_pt)
                    p.first_line_indent = Pt(-18)
                if not bullet_added:
                    if is_ordered:
                        counters[counter_key] += 1
                        if list_style_type == 'decimal':
                            marker_str = f"{counters[counter_key]}."
                        elif list_style_type == 'lower-alpha':
                            marker_str = f"{chr(96 + counters[counter_key])}."
                        elif list_style_type == 'upper-alpha':
                            marker_str = f"{chr(64 + counters[counter_key])}."
                        else:
                            marker_str = f"{counters[counter_key]}."
                        marker_run = p.add_run()
                        marker_run.text = marker_str + ' '
                        marker_run.font.name = default_font_name
                        marker_run.font.size = Pt(default_font_size_pt)
                        if default_color:
                            marker_run.font.color.rgb = default_color
                    else:
                        bullet_char = {
                            'disc': '\u2022',
                            'circle': '\u25e6',
                            'square': '\u25aa'
                        }.get(list_style_type, '\u2022')
                        bullet_run = p.add_run()
                        bullet_run.text = bullet_char + ' '
                        bullet_run.font.name = 'Arial'
                        bullet_run.font.size = Pt(default_font_size_pt)
                        if default_color:
                            bullet_run.font.color.rgb = default_color
                    bullet_added = True
                if first:
                    element_text = element_text.lstrip()
                if not element_text.strip():
                    continue
                first = False
                run = p.add_run()
                run.text = element_text
                inline_styles = inline_element.get('styles', {})
                font = run.font
                font_size_px = safe_float(inline_styles.get('fontSize', '16').replace('px', ''))
                font.name = inline_styles.get('fontFamily', 'Segoe UI').split(',')[0].strip('"\'')
                font.size = Pt(get_font_size_pt(font_size_px))
                font.bold = inline_styles.get('fontWeight', '400') in ['bold', '600', '700', '800', '900']
                font.italic = inline_styles.get('fontStyle', 'normal') == 'italic'
                color = parse_color(inline_styles.get('color'))
                if color:
                    font.color.rgb = color
            # Trim trailing spaces from the last run in the last paragraph
            if text_frame.paragraphs and text_frame.paragraphs[-1].runs:
                last_run = text_frame.paragraphs[-1].runs[-1]
                last_run.text = last_run.text.rstrip()
        else:
            if first_item and level == 0:
                p = text_frame.paragraphs[0]  # Use first paragraph for first item
            else:
                p = text_frame.add_paragraph()
            p.level = level
            p.space_after = Pt(space_after_pt)
            p.line_spacing = line_spacing
            indent_pt = 18 * level
            p.left_indent = Pt(indent_pt)
            p.first_line_indent = Pt(-18)
            if is_ordered:
                counters[counter_key] += 1
                if list_style_type == 'decimal':
                    marker_str = f"{counters[counter_key]}."
                elif list_style_type == 'lower-alpha':
                    marker_str = f"{chr(96 + counters[counter_key])}."
                elif list_style_type == 'upper-alpha':
                    marker_str = f"{chr(64 + counters[counter_key])}."
                else:
                    marker_str = f"{counters[counter_key]}."
                marker_run = p.add_run()
                marker_run.text = marker_str + ' '
                marker_run.font.name = default_font_name
                marker_run.font.size = Pt(default_font_size_pt)
                if default_color:
                    marker_run.font.color.rgb = default_color
            else:
                bullet_char = {
                    'disc': '\u2022',
                    'circle': '\u25e6',
                    'square': '\u25aa'
                }.get(list_style_type, '\u2022')
                bullet_run = p.add_run()
                bullet_run.text = bullet_char + ' '
                bullet_run.font.name = 'Arial'
                bullet_run.font.size = Pt(default_font_size_pt)
                if default_color:
                    bullet_run.font.color.rgb = default_color
            run = p.add_run()
            run.text = item.get('text', '').strip()
            font = run.font
            font.name = default_font_name
            font.size = Pt(default_font_size_pt)
            font.bold = item_styles.get('fontWeight', '400') in ['bold', '600', '700', '800', '900']
            font.italic = item_styles.get('fontStyle', 'normal') == 'italic'
            color = parse_color(item_styles.get('color'))
            if color:
                font.color.rgb = color
        
        first_item = False  # Set to False after first item
        if item.get('nestedList'):
            add_list_paragraphs(text_frame, item['nestedList'], level + 1, counters)

def add_list_element(slide, element, slide_width, slide_height, parent_has_shadow=False):
    list_info = element.get('listInfo', {})
    if not list_info.get('items'):
        return
    rect = list_info.get('rect', {})
    x = safe_int(rect.get('x', 0))
    y = safe_int(rect.get('y', 0))
    width = safe_int(rect.get('width', element.get('width', 100)))
    height = safe_int(rect.get('height', element.get('height', 100)))
    styles = element.get('styles', {})
    box_shadow = styles.get('boxShadow', 'none')
    has_shadow = box_shadow != 'none' and not parent_has_shadow
    bg_color = parse_color(styles.get('backgroundColor'))
    border_radius_str = styles.get('borderRadius', '0px')
    border_radius = parse_border_radius(border_radius_str, width, height)
    has_radius = border_radius > 0
    has_border = is_uniform_border(styles)
    has_any_border_sides = has_any_border(styles)
    try:
        if bg_color or has_border or has_any_border_sides or has_radius or has_shadow:
            add_bg_shape(slide, styles, x, y, width, height)
        textbox = slide.shapes.add_textbox(
            pixels_to_emu(x), pixels_to_emu(y),
            pixels_to_emu(width), pixels_to_emu(height)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.margin_left = pixels_to_emu(safe_float(styles.get('paddingLeft', '0px').replace('px', '')))
        text_frame.margin_right = pixels_to_emu(safe_float(styles.get('paddingRight', '0px').replace('px', '')))
        text_frame.margin_top = pixels_to_emu(safe_float(styles.get('paddingTop', '0px').replace('px', '')))
        text_frame.margin_bottom = pixels_to_emu(safe_float(styles.get('paddingBottom', '0px').replace('px', '')))
        textbox.fill.background()
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        text_frame.clear()
        add_list_paragraphs(text_frame, list_info)
    except Exception as e:
        print(f"Failed to add list: {e}")

def set_cell_border(cell, side, width_pt, color_rgb, style='solid'):
    if width_pt <= 0 or not color_rgb:
        return
    side_tag = side.lower()
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = tcPr.find(qn('a:tcBorders'))
    if tcBorders is None:
        tcBorders = OxmlElement('a:tcBorders')
        tcPr.append(tcBorders)
    existing_side = tcBorders.find(qn(f'a:{side_tag}'))
    if existing_side is not None:
        tcBorders.remove(existing_side)
    side_elem = OxmlElement(f'a:{side_tag}')
    ln = OxmlElement('a:ln')
    ln.set('w', str(int(width_pt * 12700)))
    solidFill = OxmlElement('a:solidFill')
    srgbClr = OxmlElement('a:srgbClr')
    if color_rgb is not None and hasattr(color_rgb, 'r') and hasattr(color_rgb, 'g') and hasattr(color_rgb, 'b'):
        srgbClr.set('val', f'{color_rgb.r:02X}{color_rgb.g:02X}{color_rgb.b:02X}')
    else:
        return
    solidFill.append(srgbClr)
    ln.append(solidFill)
    if style != 'solid':
        prstDash = OxmlElement('a:prstDash')
        val = 'dash' if style == 'dashed' else 'sysDot' if style == 'dotted' else 'solid'
        prstDash.set('val', val)
        ln.append(prstDash)
    side_elem.append(ln)
    tcBorders.append(side_elem)

def add_table_element(slide, element, slide_width, slide_height, parent_has_shadow=False):
    table_info = element.get('tableInfo', {})
    if not table_info.get('rows'):
        return
    rect = table_info.get('rect', {})
    x = safe_int(rect.get('x', element.get('x', 0)))
    y = safe_int(rect.get('y', element.get('y', 0)))
    width = safe_int(rect.get('width', element.get('width', 100)))
    height = safe_int(rect.get('height', element.get('height', 100)))
    styles = table_info.get('styles', {})
    box_shadow = styles.get('boxShadow', 'none')
    has_shadow = box_shadow != 'none' and not parent_has_shadow
    bg_color = parse_color(styles.get('backgroundColor'))
    border_radius_str = styles.get('borderRadius', '0px')
    border_radius = parse_border_radius(border_radius_str, width, height)
    has_radius = border_radius > 0
    has_border = is_uniform_border(styles)
    has_any_border_sides = has_any_border(styles)
    try:
        if bg_color or has_border or has_any_border_sides or has_radius or has_shadow:
            add_bg_shape(slide, styles, x, y, width, height)
        rows = table_info['rowCount']
        cols = table_info['columnCount']
        table_shape = slide.shapes.add_table(
            rows, cols, 
            pixels_to_emu(x), pixels_to_emu(y), 
            pixels_to_emu(width), pixels_to_emu(height))
        table = table_shape.table
        # Column widths
        total_html_width = 0
        col_html_widths = [0] * cols
        for row_data in table_info['rows']:
            cell_idx = 0
            for cell_data in row_data['cells']:
                cell_width = cell_data.get('rect', {}).get('width', 0)
                col_span = cell_data.get('colSpan', 1)
                for _ in range(col_span):
                    col_html_widths[cell_idx] = max(col_html_widths[cell_idx], cell_width / col_span)
                    cell_idx += 1
        total_html_width = sum(col_html_widths)
        if total_html_width > 0:
            for col_idx in range(cols):
                proportional_width = (col_html_widths[col_idx] / total_html_width) * width
                table.columns[col_idx].width = pixels_to_emu(max(10, proportional_width))
        for row_idx, row_data in enumerate(table_info['rows']):
            table.rows[row_idx].height = pixels_to_emu(safe_float(row_data.get('rect').get('height', height / rows)))
        for row_data in table_info['rows']:
            row_index = row_data['index']
            if row_index >= rows:
                continue
            row_bg_color = parse_color(row_data.get('styles', {}).get('backgroundColor'))
            for cell_data in row_data['cells']:
                cell_index = cell_data['cellIndex']
                if cell_index >= cols:
                    continue
                pptx_cell = table.cell(row_index, cell_index)
                col_span = cell_data.get('colSpan', 1)
                row_span = cell_data.get('rowSpan', 1)
                if col_span > 1 or row_span > 1:
                    try:
                        end_row = min(rows - 1, row_index + row_span - 1)
                        end_col = min(cols - 1, cell_index + col_span - 1)
                        pptx_cell.merge(table.cell(end_row, end_col))
                    except:
                        pass
                text_frame = pptx_cell.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell_styles = cell_data.get('styles', {})
                text_frame.margin_left = pixels_to_emu(safe_float(cell_styles.get('paddingLeft', '8px').replace('px', '')))
                text_frame.margin_right = pixels_to_emu(safe_float(cell_styles.get('paddingRight', '8px').replace('px', '')))
                text_frame.margin_top = pixels_to_emu(safe_float(cell_styles.get('paddingTop', '8px').replace('px', '')))
                text_frame.margin_bottom = pixels_to_emu(safe_float(cell_styles.get('paddingBottom', '8px').replace('px', '')))
                bg_color_cell = parse_color(cell_styles.get('backgroundColor'))
                if bg_color_cell:
                    pptx_cell.fill.solid()
                    pptx_cell.fill.fore_color.rgb = bg_color_cell
                elif row_bg_color:
                    pptx_cell.fill.solid()
                    pptx_cell.fill.fore_color.rgb = row_bg_color
                else:
                    pptx_cell.fill.background()
                sides = ['left', 'right', 'top', 'bottom']
                for low_side in sides:
                    cap_side = low_side.capitalize()
                    width_str = cell_styles.get(f'border{cap_side}Width', '0px')
                    width_pt = safe_float(width_str.replace('px', ''))
                    color = parse_color(cell_styles.get(f'border{cap_side}Color'))
                    style = cell_styles.get(f'border{cap_side}Style', 'solid')
                    set_cell_border(pptx_cell, low_side, width_pt, color, style)
                text_frame.clear()
                p = text_frame.paragraphs[0]  # Use the first paragraph instead of adding new one
                first = True
                text_align = cell_styles.get('textAlign', 'left')
                alignment = PP_ALIGN.CENTER if text_align == 'center' else PP_ALIGN.RIGHT if text_align == 'right' else PP_ALIGN.LEFT
                p.alignment = alignment
                if cell_data.get('inlineGroup') and cell_data['inlineGroup'].get('inlineElements'):
                    for inline_element in cell_data['inlineGroup']['inlineElements']:
                        if inline_element.get('type') == 'br':
                            if p is not None:
                                p = text_frame.add_paragraph()
                                p.alignment = alignment
                            first = True
                            continue
                        element_text = inline_element.get('text', '')
                        if first:
                            element_text = element_text.lstrip()
                        if not element_text.strip():
                            continue
                        first = False
                        run = p.add_run()
                        run.text = element_text
                        inline_styles = inline_element.get('styles', {})
                        font = run.font
                        font_size_px = safe_float(inline_styles.get('fontSize', '14').replace('px', ''))
                        font.name = inline_styles.get('fontFamily', 'Segoe UI').split(',')[0].strip('"\'')
                        font.size = Pt(max(9, get_font_size_pt(font_size_px)))
                        font.bold = inline_styles.get('fontWeight', '400') in ['bold', '600', '700', '800', '900']
                        font.italic = inline_styles.get('fontStyle', 'normal') == 'italic'
                        color = parse_color(inline_styles.get('color'))
                        if color:
                            font.color.rgb = color
                    # Trim trailing spaces from the last run in the last paragraph
                    if text_frame.paragraphs and text_frame.paragraphs[-1].runs:
                        last_run = text_frame.paragraphs[-1].runs[-1]
                        last_run.text = last_run.text.rstrip()
                else:
                    run = p.add_run()
                    run.text = cell_data.get('text', '').strip()
                    font = run.font
                    font_size_px = safe_float(cell_styles.get('fontSize', '14').replace('px', ''))
                    font.name = cell_styles.get('fontFamily', 'Segoe UI').split(',')[0].strip('"\'')
                    font.size = Pt(max(9, get_font_size_pt(font_size_px)))
                    font.bold = cell_styles.get('fontWeight', '400') in ['bold', '600', '700', '800', '900']
                    font.italic = cell_styles.get('fontStyle', 'normal') == 'italic'
                    color = parse_color(cell_styles.get('color'))
                    if color:
                        font.color.rgb = color
    except Exception as e:
        print(f"Failed to add table: {e}")

def add_image_element(slide, element, slide_width, slide_height, parent_has_shadow=False):
    media_info = element.get('mediaInfo', {})
    img_src = media_info.get('src', '')
    styles = element.get('styles', {})
    if not img_src:
        return
    # Use precise positioning from extraction
    x = element.get('x', 0)
    y = element.get('y', 0)
    width = max(1, element.get('width', 100))
    height = max(1, element.get('height', 100))
    
    # Ensure coordinates are within slide bounds
    x = max(0, min(x, slide_width - width))
    y = max(0, min(y, slide_height - height))
    
    natural_width = media_info.get('naturalWidth', width)
    border_radius_str = styles.get('borderRadius', '0px')
    radius_ratio = parse_border_radius(border_radius_str, width, height)
    radius_display = radius_ratio * min(width, height)
    has_radius = radius_display > 0
    box_shadow = styles.get('boxShadow', 'none')
    has_shadow = box_shadow != 'none' and not parent_has_shadow
    has_border = is_uniform_border(styles)
    
    try:
        temp_path = None
        if img_src.startswith('data:'):
            _, data = img_src.split(',', 1)
            img_data = base64.b64decode(data)
            temp_path = 'temp_image.png'
            with open(temp_path, 'wb') as f:
                f.write(img_data)
        elif img_src.startswith('http'):
            response = requests.get(img_src, timeout=10)
            if response.status_code == 200:
                temp_path = 'temp_image.png'
                with open(temp_path, 'wb') as f:
                    f.write(response.content)
            else:
                print(f"Failed to download image: {img_src}")
                return
        else:
            if os.path.exists(img_src):
                temp_path = img_src
            else:
                print(f"Image file not found: {img_src}")
                return
        
        with Image.open(temp_path) as img:
            img.verify()
        
        image_to_add = temp_path
        if has_radius:
            scale_x = natural_width / width if width > 0 else 1
            radius_natural = int(radius_display * scale_x)
            temp_rounded = 'temp_rounded.png'
            make_rounded_image(temp_path, temp_rounded, radius_natural)
            image_to_add = temp_rounded
        
        # Add image with precise positioning
        picture = slide.shapes.add_picture(
            image_to_add,
            pixels_to_emu(x), pixels_to_emu(y),
            pixels_to_emu(width), pixels_to_emu(height)
        )
        picture.shadow.inherit = False
        
        # Handle borders and shadows
        if has_border or has_radius or has_shadow:
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if has_radius else MSO_SHAPE.RECTANGLE
            border_shape = slide.shapes.add_shape(
                shape_type,
                pixels_to_emu(x), pixels_to_emu(y),
                pixels_to_emu(width), pixels_to_emu(height)
            )
            if has_radius:
                border_shape.adjustments[0] = radius_ratio
            border_shape.fill.background()
            border_shape.shadow.inherit = False
            if has_shadow:
                apply_shadow(border_shape, box_shadow)
            if has_border:
                border_width = safe_float(styles.get('borderTopWidth', '0px'))
                border_color = parse_color(styles.get('borderTopColor'))
                if border_color:
                    border_shape.line.width = Pt(border_width)
                    border_shape.line.color.rgb = border_color
                    border_style = styles.get('borderTopStyle', 'solid')
                    if border_style == 'dashed':
                        border_shape.line.dash_style = MSO_LINE.DASH
                    elif border_style == 'dotted':
                        border_shape.line.dash_style = MSO_LINE.ROUND_DOT
            else:
                border_shape.line.fill.background()
            # Move border shape behind picture
            sp = border_shape._sp
            parent = sp.getparent()
            parent.remove(sp)
            pic_sp = picture._sp
            idx = list(parent).index(pic_sp)
            parent.insert(idx, sp)
        elif has_shadow:
            apply_shadow(picture, box_shadow)
        
        # Clean up temporary files
        if image_to_add != temp_path and os.path.exists(image_to_add):
            os.remove(image_to_add)
        if temp_path != img_src and os.path.exists(temp_path):
            os.remove(temp_path)
    except Exception as e:
        print(f"Failed to add image: {e}")

def add_text_element(slide, element, slide_width, slide_height, parent_has_shadow=False):
    """Enhanced text element creation with precise positioning"""
    text = element.get('text', '').strip()
    if not text:
        return
    
    # Use precise positioning
    x = element.get('x', 0)
    y = element.get('y', 0)
    width = max(1, element.get('width', 100))
    height = max(1, element.get('height', 100))
    
    # Ensure coordinates are within slide bounds
    x = max(0, min(x, slide_width - width))
    y = max(0, min(y, slide_height - height))
    
    styles = element.get('styles', {})
    box_shadow = styles.get('boxShadow', 'none')
    has_shadow = box_shadow != 'none' and not parent_has_shadow
    bg_color = parse_color(styles.get('backgroundColor'))
    border_radius_str = styles.get('borderRadius', '0px')
    border_radius = parse_border_radius(border_radius_str, width, height)
    has_radius = border_radius > 0
    has_border = is_uniform_border(styles)
    has_any_border_sides = has_any_border(styles)
    
    try:
        # Add background/border shapes first
        if bg_color or has_border or has_any_border_sides or has_radius or has_shadow:
            add_bg_shape(slide, styles, x, y, width, height)
        
        # Create text box with precise positioning
        textbox = slide.shapes.add_textbox(
            pixels_to_emu(x), pixels_to_emu(y),
            pixels_to_emu(width), pixels_to_emu(height)
        )
        
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # Set vertical alignment based on flex properties or default to middle for short text
        display = styles.get('display', 'block')
        align_items = styles.get('alignItems', 'stretch')
        justify_content = styles.get('justifyContent', 'flex-start')
        
        if display == 'flex' and align_items == 'center':
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        elif height < 50:  # For small elements like company names
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        # Apply precise margins
        text_frame.margin_left = pixels_to_emu(safe_float(styles.get('paddingLeft', '0px')))
        text_frame.margin_right = pixels_to_emu(safe_float(styles.get('paddingRight', '0px')))
        text_frame.margin_top = pixels_to_emu(safe_float(styles.get('paddingTop', '0px')))
        text_frame.margin_bottom = pixels_to_emu(safe_float(styles.get('paddingBottom', '0px')))
        
        # Remove textbox styling
        textbox.fill.background()
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        
        # Clear and set text
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        
        # Apply text formatting
        font = run.font
        font_size_px = safe_float(styles.get('fontSize', '12'))
        font.name = styles.get('fontFamily', 'Arial').split(',')[0].strip('"\'')
        font.size = Pt(max(6, get_font_size_pt(font_size_px)))
        font.bold = styles.get('fontWeight', '400') in ['bold', '700', '800', '900']
        font.italic = styles.get('fontStyle') == 'italic'
        
        # Apply text color
        color = parse_color(styles.get('color', 'black'))
        if color:
            font.color.rgb = color
        
        # Apply text alignment based on CSS text-align and flex properties
        text_align = styles.get('textAlign', 'left')
        if text_align == 'center' or (display == 'flex' and justify_content == 'center'):
            p.alignment = PP_ALIGN.CENTER
        elif text_align == 'right' or (display == 'flex' and justify_content == 'flex-end'):
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
            
    except Exception as e:
        print(f"Failed to add text: {e}")

def get_parent(element, all_elements):
    el_x = element.get('x', 0)
    el_y = element.get('y', 0)
    el_w = element.get('width', 0)
    el_h = element.get('height', 0)
    el_rect = (el_x, el_y, el_x + el_w, el_y + el_h)
    potential_parents = []
    for other in all_elements:
        if other is element:
            continue
        o_x = other.get('x', 0)
        o_y = other.get('y', 0)
        o_w = other.get('width', 0)
        o_h = other.get('height', 0)
        o_rect = (o_x, o_y, o_x + o_w, o_y + o_h)
        if el_rect[0] >= o_rect[0] and el_rect[1] >= o_rect[1] and el_rect[2] <= o_rect[2] and el_rect[3] <= o_rect[3]:
            area = o_w * o_h
            potential_parents.append((area, other))
    if potential_parents:
        potential_parents.sort(key=lambda p: p[0])  # smallest area first
        return potential_parents[0][1]
    return None

def create_pptx_from_json(json_path, output_path=None):
    """Enhanced PowerPoint generation with precise positioning"""
    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
    except Exception as e:
        print(f"Error reading JSON file: {e}")
        return
    
    if not slides_data:
        print("No slides found in JSON")
        return
    
    # Get slide dimensions from first slide
    first_slide = slides_data[0]
    slide_width = safe_int(first_slide.get('slideWidth', 1920))
    slide_height = safe_int(first_slide.get('slideHeight', 1080))
    
    # Create presentation with precise dimensions
    prs = Presentation()
    prs.slide_width = pixels_to_emu(slide_width)
    prs.slide_height = pixels_to_emu(slide_height)
    
    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = parse_color('#ffffff') or RGBColor(255, 255, 255)
        
        # Add slide background styling
        slide_styles = slide_data.get('slideStyles', {})
        if slide_styles:
            add_bg_shape(slide, slide_styles, 0, 0, slide_width, slide_height)
        
        elements = slide_data.get('elements', [])
        
        # Enhanced sorting: separate background elements from content elements
        # Background elements (divs without content) should render first
        # Images and text should render last to stay on top
        def get_element_priority(element):
            element_type = element.get('type', '').lower()
            has_text = bool(element.get('text', '').strip())
            has_inline_group = bool(element.get('inlineGroup'))
            has_image = element_type == 'img'
            has_table = element_type == 'table'
            has_list = element_type in ['ul', 'ol']
            
            # Priority order (lower number = rendered first/behind)
            if element_type == 'div' and not has_text and not has_inline_group:
                return (0, element.get('zIndex', 0), element.get('y', 0), element.get('x', 0))  # Background divs first
            elif has_list or has_table:
                return (1, element.get('zIndex', 0), element.get('y', 0), element.get('x', 0))  # Lists and tables
            elif has_text or has_inline_group:
                return (2, element.get('zIndex', 0), element.get('y', 0), element.get('x', 0))  # Text elements
            elif has_image:
                return (3, element.get('zIndex', 0), element.get('y', 0), element.get('x', 0))  # Images on top
            else:
                return (1, element.get('zIndex', 0), element.get('y', 0), element.get('x', 0))  # Other elements
        
        elements_sorted = sorted(elements, key=get_element_priority)
        
        # Build parent hierarchy for shadow inheritance
        parent_map = {}
        for el in elements_sorted:
            parent = get_parent(el, elements_sorted)
            parent_map[id(el)] = parent
        
        # Process each element with enhanced positioning
        for element in elements_sorted:
            element_type = element.get('type', '').lower()
            
            # --- Enhancement: handle .company and .footer children as separate elements ---
            if element_type == 'div' and 'company' in element.get('className', ''):
                # Find and render background first, then children
                # Render background div if it has styling
                styles = element.get('styles', {})
                if (has_any_border(styles) or 
                    parse_color(styles.get('backgroundColor')) or
                    styles.get('boxShadow', 'none') != 'none'):
                    x = element.get('x', 0)
                    y = element.get('y', 0)
                    width = max(1, element.get('width', 100))
                    height = max(1, element.get('height', 100))
                    add_bg_shape(slide, styles, x, y, width, height)
                
                # Then render children on top
                for child in elements_sorted:
                    child_x = child.get('x', 0)
                    child_y = child.get('y', 0)
                    elem_x = element.get('x', 0)
                    elem_y = element.get('y', 0)
                    
                    # Check if child is within this element's bounds
                    if (child_x >= elem_x and child_y >= elem_y and
                        child_x < elem_x + element.get('width', 0) and
                        child_y < elem_y + element.get('height', 0)):
                        
                        if child.get('type') == 'img':
                            add_image_element(slide, child, slide_width, slide_height)
                        elif child.get('type') == 'span':
                            add_text_element(slide, child, slide_width, slide_height)
                continue
                
            if element_type == 'div' and 'footer' in element.get('className', ''):
                # Render background first
                styles = element.get('styles', {})
                if (has_any_border(styles) or 
                    parse_color(styles.get('backgroundColor')) or
                    styles.get('boxShadow', 'none') != 'none'):
                    x = element.get('x', 0)
                    y = element.get('y', 0)
                    width = max(1, element.get('width', 100))
                    height = max(1, element.get('height', 100))
                    add_bg_shape(slide, styles, x, y, width, height)
                
                # Then render children on top
                elem_y = element.get('y', 0)
                for child in elements_sorted:
                    if child.get('y', 0) >= elem_y:
                        if child.get('type') == 'img':
                            add_image_element(slide, child, slide_width, slide_height)
                        elif child.get('type') == 'span':
                            add_text_element(slide, child, slide_width, slide_height)
                continue

            if element_type == 'canvas':
                continue
                
            # Skip child elements of company/footer divs as they're handled above
            parent = parent_map.get(id(element))
            if parent and ('company' in parent.get('className', '') or 'footer' in parent.get('className', '')):
                continue
            
            parent_has_shadow = bool(parent and parent.get('styles', {}).get('boxShadow', 'none') != 'none')
            
            if element.get('inlineGroup'):
                add_inline_group_element(slide, element, slide_width, slide_height, parent_has_shadow)
            elif element_type in ['ul', 'ol']:
                add_list_element(slide, element, slide_width, slide_height, parent_has_shadow)
            elif element_type == 'table':
                add_table_element(slide, element, slide_width, slide_height, parent_has_shadow)
            elif element_type == 'img':
                add_image_element(slide, element, slide_width, slide_height, parent_has_shadow)
            elif element_type == 'span':
                add_text_element(slide, element, slide_width, slide_height, parent_has_shadow)
            elif element_type in ['div', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                if (element.get('text', '').strip() or
                    has_any_border(element.get('styles', {})) or
                    parse_color(element.get('styles', {}).get('backgroundColor')) or
                    element.get('styles', {}).get('boxShadow', 'none') != 'none'):
                    if element.get('text', '').strip() and not element.get('inlineGroup'):
                        add_text_element(slide, element, slide_width, slide_height, parent_has_shadow)
                    elif not element.get('text', '').strip():
                        x = element.get('x', 0)
                        y = element.get('y', 0)
                        width = max(1, element.get('width', 100))
                        height = max(1, element.get('height', 100))
                        add_bg_shape(slide, element.get('styles', {}), x, y, width, height)
    
    if output_path is None:
        base_name = os.path.splitext(os.path.basename(json_path))[0]
        output_path = f"{base_name}_output.pptx"
    try:
        prs.save(output_path)
        print(f"Presentation saved as '{output_path}' with {len(slides_data)} slide(s)")
        print(f"Slide dimensions: {slide_width}x{slide_height} pixels")
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == "__main__":
    create_pptx_from_json('slides_data.json', 'output.pptx')